"""Generate Dormify project presentation (PPTX)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "Dormify-Presentation.pptx"

PRIMARY = RGBColor(0x0D, 0x6E, 0xFD)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x19, 0x87, 0x54)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(1.35),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(8.8), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(0.5))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = MUTED
        sp.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items: list[str], top: float = 1.9) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.6), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(14)
        p.alignment = PP_ALIGN.RIGHT
        p.line_spacing = 1.25


def add_two_columns(slide, left_title: str, left_items: list[str], right_title: str, right_items: list[str]) -> None:
    for col, title, items, x in [
        (0, right_title, right_items, 5.15),
        (1, left_title, left_items, 0.55),
    ]:
        hdr = slide.shapes.add_textbox(Inches(x), Inches(1.75), Inches(4.2), Inches(0.45))
        hp = hdr.text_frame.paragraphs[0]
        hp.text = title
        hp.font.size = Pt(22)
        hp.font.bold = True
        hp.font.color.rgb = PRIMARY
        hp.alignment = PP_ALIGN.RIGHT if col == 0 else PP_ALIGN.LEFT

        box = slide.shapes.add_textbox(Inches(x), Inches(2.25), Inches(4.2), Inches(4.5))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(17)
            p.font.color.rgb = DARK
            p.space_after = Pt(10)
            p.alignment = PP_ALIGN.RIGHT if col == 0 else PP_ALIGN.LEFT


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(10), Inches(1.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = DARK
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.2))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Dormify"
    tp.font.size = Pt(54)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.4), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "منصة ذكية لإسكان الطلاب الجامعيين"
    sp.font.size = Pt(26)
    sp.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    sp.alignment = PP_ALIGN.CENTER

    tag = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.6))
    tgp = tag.text_frame.paragraphs[0]
    tgp.text = "React + Express + MySQL | عرض مشروع"
    tgp.font.size = Pt(16)
    tgp.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    tgp.alignment = PP_ALIGN.CENTER

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(8.4), Inches(0.5))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "2026"
    fp.font.size = Pt(14)
    fp.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    fp.alignment = PP_ALIGN.CENTER


def slide_content(prs: Presentation, title: str, subtitle: str | None, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, title, subtitle)
    add_bullets(slide, bullets)


def slide_two_col(prs: Presentation, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, title)
    add_two_columns(slide, left_title, left_items, right_title, right_items)


def slide_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "سير العمل", "من التسجيل إلى تأكيد الحجز")

    steps = [
        ("1", "تسجيل / دخول", "طالب أو مالك"),
        ("2", "تصفح السكن", "بحث وفلترة"),
        ("3", "طلب حجز", "حالة: قيد الانتظار"),
        ("4", "مراجعة المالك", "تأكيد أو رفض"),
        ("5", "تحديث الغرف", "تنقص available_rooms"),
    ]
    y_start = 2.0
    for i, (num, label, desc) in enumerate(steps):
        y = y_start + i * 0.95
        circle = slide.shapes.add_shape(9, Inches(8.5), Inches(y), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        nbox = slide.shapes.add_textbox(Inches(8.5), Inches(y + 0.05), Inches(0.55), Inches(0.45))
        np = nbox.text_frame.paragraphs[0]
        np.text = num
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.alignment = PP_ALIGN.CENTER

        lbl = slide.shapes.add_textbox(Inches(0.7), Inches(y - 0.05), Inches(7.5), Inches(0.4))
        lp = lbl.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(20)
        lp.font.bold = True
        lp.font.color.rgb = DARK
        lp.alignment = PP_ALIGN.RIGHT

        dbox = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(7.5), Inches(0.35))
        dp = dbox.text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(15)
        dp.font.color.rgb = MUTED
        dp.alignment = PP_ALIGN.RIGHT


def slide_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.4), Inches(1))
    tp = t.text_frame.paragraphs[0]
    tp.text = "شكراً لكم"
    tp.font.size = Pt(44)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER

    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.8))
    sp = s.text_frame.paragraphs[0]
    sp.text = "Dormify — Smart Student Housing"
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    sp.alignment = PP_ALIGN.CENTER

    c = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(8.4), Inches(0.6))
    cp = c.text_frame.paragraphs[0]
    cp.text = "info@dormify.com  |  github.com/Ragheed173/Dormify-Website"
    cp.font.size = Pt(14)
    cp.font.color.rgb = PRIMARY
    cp.alignment = PP_ALIGN.CENTER


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)

    slide_content(
        prs,
        "المشكلة",
        "تحديات الطلاب في البحث عن سكن",
        [
            "صعوبة إيجاد سكن قريب من الجامعة بسعر مناسب",
            "قوائم متفرقة على مجموعات ووسائل تواصل مختلفة",
            "غياب نظام موحد للحجز ومتابعة الطلبات",
            "قلق الملاك من إدارة الطلبات والغرف المتاحة يدوياً",
        ],
    )

    slide_content(
        prs,
        "الحل — Dormify",
        "منصة ويب متكاملة لربط الطلاب بملاك السكن",
        [
            "منصة واحدة لتصفح السكن، الحجز، وإدارة الطلبات",
            "ثلاثة أدوار: طالب، مالك سكن، ومشرف (Admin)",
            "واجهة React حديثة مع لوحات تحكم لكل دور",
            "خادم Express مع قاعدة بيانات MySQL عبر Sequelize",
        ],
    )

    slide_two_col(
        prs,
        "المستخدمون والأدوار",
        "الطالب",
        [
            "تسجيل ودخول (JWT / Google)",
            "تصفح السكن مع بحث وفلترة",
            "إنشاء طلبات حجز",
            "متابعة حالة الحجوزات",
        ],
        "المالك والمشرف",
        [
            "إضافة وتعديل إعلانات السكن",
            "تأكيد / رفض / إلغاء الحجوزات",
            "إدارة المستخدمين والإحصائيات (Admin)",
            "مراقبة المنصة بالكامل",
        ],
    )

    slide_content(
        prs,
        "ميزات المنصة",
        None,
        [
            "قوائم سكن عامة مع pagination وبحث متقدم",
            "نظام حجز: Pending → Confirmed / Rejected / Cancelled",
            "تحديث تلقائي لعدد الغرف المتاحة عند التأكيد",
            "رفع صور للسكن (Multer) وإشعارات بريد (Nodemailer)",
            "توثيق API عبر Swagger على /api-docs",
        ],
    )

    slide_two_col(
        prs,
        "المساعد الذكي (AI)",
        "شرح المفاهيم",
        [
            "POST /api/ai/explain",
            "شرح مواضيع تقنية للطلاب",
            "دعم Gemini أو Groq",
        ],
        "بحث السكن بالذكاء الاصطناعي",
        [
            "POST /api/ai/housing-search",
            "طلب بلغة طبيعية: «غرفة قرب الجامعة»",
            "مطابقة مع قوائم Dormify الفعلية",
            "وضع تجريبي: USE_MOCK_AI=1",
        ],
    )

    slide_content(
        prs,
        "الأمان والبنية",
        None,
        [
            "مصادقة JWT + تسجيل دخول Google OAuth",
            "Middleware للأدوار: student / owner / admin",
            "التحقق من الطلبات (validateRequest) ومعالجة أخطاء مركزية",
            "هيكل: Routes → Controllers → Services → Models",
            "bookingService.js لمنطق حالة الحجز المشترك",
        ],
    )

    slide_two_col(
        prs,
        "التقنيات المستخدمة",
        "Frontend",
        [
            "React 18 + Vite",
            "React Router",
            "Bootstrap 5 + Icons",
            "Context API للمصادقة",
        ],
        "Backend",
        [
            "Node.js + Express 5",
            "MySQL + Sequelize",
            "Passport (Google OAuth)",
            "Swagger, bcrypt, JWT",
        ],
    )

    slide_flow(prs)

    slide_content(
        prs,
        "تشغيل المشروع محلياً",
        None,
        [
            "نسخ .env.example إلى .env وتعبئة المتغيرات",
            "npm install ثم npm run backend (منفذ 5000)",
            "npm run frontend في طرفية ثانية (منفذ 5173)",
            "Swagger: http://localhost:5000/api-docs",
            "تطوير DB: DB_SYNC_ALTER=true عند تغيير المخطط",
        ],
    )

    slide_closing(prs)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
